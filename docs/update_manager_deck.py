#!/usr/bin/env python3
"""
Refresh manager deck with POC findings (no SCR/CAL step-level detail).

Warning: run once per deck version — re-running adds duplicate slides.
After edits, slide order should be:
  1 Title → 2 Scope → 3 Architecture → 4 Capabilities → 5 Findings
  → 6 Business Outputs → 7 vs Matterway → 8 Technology → 9 Deployment
  → 10 Future → 11 Conclusion
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

DECK_PATH = Path(__file__).resolve().parent / "Document_Processing_POC_Manager_Deck_7_Slides.pptx"

MARGIN_LEFT = 731520
BULLET_TOP_START = 1005840
BULLET_STEP = 520000
BULLET_HEIGHT = 580000
CONTENT_WIDTH = 10515600
FONT_BODY = Pt(20)
FONT_TITLE = Pt(28)


def _set_title(slide, title: str) -> None:
    shape = slide.shapes[0]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = FONT_TITLE


def _set_bullet_shapes(slide, bullets: list[str], start_index: int = 1) -> None:
    """Fill text boxes from start_index with bullets (one per shape)."""
    for i, text in enumerate(bullets):
        idx = start_index + i
        if idx >= len(slide.shapes):
            break
        sh = slide.shapes[idx]
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        line = text if text.startswith("•") else f"• {text}"
        p.text = line
        p.font.size = FONT_BODY


def _duplicate_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(prs.slide_layouts[6])
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return dest


def _insert_slide_at(prs: Presentation, slide_index: int) -> None:
    """Move last slide to slide_index."""
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(slide_index, new_id)


def _add_content_slide(prs: Presentation, template_index: int, title: str, bullets: list[str]) -> None:
    _duplicate_slide(prs, template_index)
    _insert_slide_at(prs, len(prs.slides) - 1)  # will be moved again by caller
    slide = prs.slides[-1]
    _set_title(slide, title)
    # Remove extra template bullets beyond our count
    while len(slide.shapes) > 1 + len(bullets):
        sp = slide.shapes[1]
        slide.shapes._spTree.remove(sp._element)
    # Add text boxes if needed
    while len(slide.shapes) < 1 + len(bullets):
        top = BULLET_TOP_START + (len(slide.shapes) - 1) * BULLET_STEP
        box = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, BULLET_HEIGHT)
        box.text_frame.word_wrap = True
    _set_bullet_shapes(slide, bullets)


def update_deck() -> None:
    prs = Presentation(str(DECK_PATH))

    # --- Slide 1: Title ---
    _set_title(prs.slides[0], "Document Processing POC — Findings")
    _set_bullet_shapes(
        prs.slides[0],
        [
            "Confirm Excel, Word & PDF automation is feasible off Matterway (owned server + web app).",
            "Demonstrate the same class of capability SCR/CAL rely on — not a full SCR/CAL replacement.",
            "De-risk hosting on Windows/IIS with domain, API, and operational logging.",
        ],
    )

    # --- Slide 2: Scope ---
    _set_title(prs.slides[1], "POC Objective & Scope")
    _set_bullet_shapes(
        prs.slides[1],
        [
            "Objective: prove server-side read/write Excel and read/validate PDF & Word at scale.",
            "In scope: upload & multi-upload, server path, downloads, auto-download option, in-place Excel, logs.",
            "Out of scope (by design): Matterway SCR scrub & CAL allocation products — logic not ported; feasible next.",
            "Production next: authentication, job queue, audit trail, environment-based rules.",
        ],
    )

    # --- Slide 3: Architecture ---
    _set_title(prs.slides[2], "Solution Architecture")
    _set_bullet_shapes(
        prs.slides[2],
        [
            "User → React UI → FastAPI → Excel / PDF / Word services → results & downloads.",
            "Modular processors; config-driven PDF/Word validation; REST API for future integrations.",
            "Upload or authorized server path; optional in-place update on configured workbook path.",
        ],
        start_index=1,
    )

    # --- Slide 4: Capabilities ---
    _set_title(prs.slides[3], "Capabilities Demonstrated")
    _set_bullet_shapes(
        prs.slides[3],
        [
            "Excel (.xlsx/.xls): POC Status column, FALSE/NA highlighting, Name & SSN worksheet, processed download.",
            "PDF & Word: text extraction; PASS/FAIL vs required headings, questions, answers (config files).",
            "Large-file testing; processing time on UI; upload copy removed after success; processed Excel retained.",
            "Optional auto-download after Excel processing; static in-place path update (OneDrive/local).",
        ],
    )

    # --- Slide 5: Outcomes ---
    _set_title(prs.slides[4], "POC Findings")
    _set_bullet_shapes(
        prs.slides[4],
        [
            "Representative large Excel, PDF, and Word files processed successfully on the server.",
            "Automation parity at platform level: Matterway skills use similar document libraries in-browser.",
            "PASS/FAIL reports and Excel metrics suitable for ops review; backend logs with timestamps.",
            "Windows Server + IIS + HTTPS deployment documented (runbook in repo).",
        ],
    )

    # --- Slide 6: Technology ---
    _set_title(prs.slides[5], "Technology Stack")
    _set_bullet_shapes(
        prs.slides[5],
        [
            "Frontend: React + Material UI (Vite build for production).",
            "Backend: FastAPI + uvicorn (Python).",
            "Excel: openpyxl + xlrd (.xls); PDF: PyMuPDF; Word: python-docx.",
            "Configuration-driven rules; standard stack for hiring, support, and cloud/on-prem hosting.",
        ],
    )

    # Insert new slides before Conclusion (index 7 when 8 slides exist)
    conclusion_index = len(prs.slides) - 1
    insert_at = conclusion_index  # insert before conclusion

    # Business outputs
    _duplicate_slide(prs, 5)
    _insert_slide_at(prs, insert_at)
    _set_title(prs.slides[insert_at], "Business Outputs")
    _set_bullet_shapes(
        prs.slides[insert_at],
        [
            "Excel: updated workbook (status column, QA highlights, Name/SSN tab) + download + optional disk update.",
            "PDF/Word: validation status, lists of missing required content, document size metrics.",
            "Operational: upload summary, processing duration, success/error on result screen, server log files.",
        ],
    )
    insert_at += 1
    conclusion_index += 1

    # Matterway comparison (high level only)
    _duplicate_slide(prs, 5)
    _insert_slide_at(prs, insert_at)
    _set_title(prs.slides[insert_at], "POC vs Matterway (SCR & CAL)")
    _set_bullet_shapes(
        prs.slides[insert_at],
        [
            "Matterway SCR: browser skill for Helpwc scrubbing — POC confirms server Excel automation; scrub product not rebuilt here.",
            "Matterway CAL: multi-file contribution allocation (Helpwc, plan docs, YEDC PDF, macro handoffs) — POC confirms PDF/Excel foundation only.",
            "POC adds: standalone app, API, server path, in-place file, domain deployment — less portal/runtime lock-in.",
            "Next phase: migrate priority SCR/CAL workflows onto this platform with UAT vs legacy outputs.",
        ],
    )
    insert_at += 1
    conclusion_index += 1

    # Future extensions (no step lists)
    _duplicate_slide(prs, 5)
    _insert_slide_at(prs, insert_at)
    _set_title(prs.slides[insert_at], "Future Extensions (Same Stack)")
    _set_bullet_shapes(
        prs.slides[insert_at],
        [
            "Batch folder processing, scheduled jobs, email/Teams notifications on PASS/FAIL.",
            "Generated PDF reports from Excel; multi-file jobs like CAL intake (wizard + single backend pipeline).",
            "Azure AD / IIS auth, retention policies, monitoring, and per-plan-year rule packs.",
        ],
    )

    # Deployment slide (was 7, index shifted)
    for i, s in enumerate(prs.slides):
        if s.shapes[0].has_text_frame and s.shapes[0].text_frame.text.startswith("Deployment"):
            dep = i
            break
    else:
        dep = 7
    _set_title(prs.slides[dep], "Deployment Plan")
    bullets_dep = [
        "Windows Server + IIS with HTTPS on company domain; React static build + FastAPI as Windows service (NSSM).",
        "IIS reverse-proxy API routes to localhost:8000; public 443 only — do not expose port 8000.",
        "DNS + SSL; storage for processed Excel; optional ALLOWED_DOCUMENT_PATH_ROOTS on Windows.",
        "Runbook: docs/Windows-Server-Deployment-Guide.docx (architecture diagram included).",
    ]
    while len(prs.slides[dep].shapes) < 1 + len(bullets_dep):
        idx = len(prs.slides[dep].shapes) - 1
        top = BULLET_TOP_START + idx * BULLET_STEP
        prs.slides[dep].shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, BULLET_HEIGHT)
    _set_bullet_shapes(prs.slides[dep], bullets_dep)

    # Conclusion (last slide)
    last = len(prs.slides) - 1
    _set_title(prs.slides[last], "Conclusion & Recommended Next Steps")
    _set_bullet_shapes(
        prs.slides[last],
        [
            "POC confirms: Excel, Word & PDF automation is possible on our owned platform (like Matterway/CAL depend on).",
            "Recommend pilot: one real workflow + production hardening (auth, monitoring, config by environment).",
            "Plan phased port of SCR/CAL business logic — platform ready; product rules are the remaining work.",
        ],
    )

    prs.save(str(DECK_PATH))
    print(f"Updated {DECK_PATH} — {len(prs.slides)} slides total.")


if __name__ == "__main__":
    update_deck()
