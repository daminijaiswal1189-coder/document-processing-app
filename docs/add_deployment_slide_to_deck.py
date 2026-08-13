#!/usr/bin/env python3
"""Add Deployment Plan slide to the manager deck (before Conclusion)."""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

DECK_PATH = Path(__file__).resolve().parent / "Document_Processing_POC_Manager_Deck_7_Slides.pptx"

MARGIN_LEFT = 731520
BULLET_TOP_START = 1005840
BULLET_STEP = 594360
BULLET_HEIGHT = 548640
CONTENT_WIDTH = 10515600
FONT_SIZE_BODY = Pt(22)
FONT_SIZE_TITLE = Pt(28)

TITLE = "Deployment Plan"

BULLETS_LEFT = [
    "Windows Server + IIS with HTTPS on your company domain.",
    "React UI: production build (frontend/dist); FastAPI (uvicorn) as a Windows service (NSSM).",
    "IIS reverse-proxies API routes to 127.0.0.1:8000 (single-domain, minimal CORS).",
]

BULLETS_RIGHT = [
    "DNS + SSL; firewall opens 443 only — keep port 8000 on localhost.",
    "Uploads/processed files and logs on server; optional server-path & in-place Excel.",
    "Runbook: docs/Windows-Server-Deployment-Guide.docx (+ .md in repo).",
]


def _duplicate_slide(prs: Presentation, source_index: int):
    """Clone slide shapes onto a new blank slide."""
    source = prs.slides[source_index]
    layout = prs.slide_layouts[6]
    dest = prs.slides.add_slide(layout)
    for shape in source.shapes:
        new_el = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return dest


def _insert_slide_before_last(prs: Presentation) -> None:
    """Move the last slide to index 6 (before Conclusion)."""
    sld_id_lst = prs.slides._sldIdLst
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(6, new_id)


def _set_shape_text(slide, shape_index: int, text: str, *, bold: bool = False, size=FONT_SIZE_BODY) -> None:
    shape = slide.shapes[shape_index]
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = bold
    p.font.size = size


def add_deployment_slide(prs: Presentation) -> None:
    # Clone Technology & Design slide layout (title bar + text boxes)
    _duplicate_slide(prs, 5)
    _insert_slide_before_last(prs)
    slide = prs.slides[6]

    _set_shape_text(slide, 0, TITLE, bold=True, size=FONT_SIZE_TITLE)

    # Shapes 1–4 on template; use two-column layout via textbox positions
    while len(slide.shapes) > 1:
        sp = slide.shapes[1]
        slide.shapes._spTree.remove(sp._element)

    half_w = CONTENT_WIDTH // 2 - 100000
    for i, text in enumerate(BULLETS_LEFT):
        top = BULLET_TOP_START + i * BULLET_STEP
        box = slide.shapes.add_textbox(MARGIN_LEFT, top, half_w, BULLET_HEIGHT + 300000)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {text}"
        p.font.size = FONT_SIZE_BODY

    right_left = MARGIN_LEFT + CONTENT_WIDTH // 2 + 100000
    for i, text in enumerate(BULLETS_RIGHT):
        top = BULLET_TOP_START + i * BULLET_STEP
        box = slide.shapes.add_textbox(right_left, top, half_w, BULLET_HEIGHT + 300000)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• {text}"
        p.font.size = FONT_SIZE_BODY


def main() -> None:
    if not DECK_PATH.is_file():
        raise FileNotFoundError(DECK_PATH)
    prs = Presentation(str(DECK_PATH))
    n_before = len(prs.slides)
    add_deployment_slide(prs)
    prs.save(str(DECK_PATH))
    print(f"Updated {DECK_PATH}: {n_before} → {len(prs.slides)} slides.")
    print("Slide 7: Deployment Plan | Slide 8: Conclusion & Next Steps")


if __name__ == "__main__":
    main()
