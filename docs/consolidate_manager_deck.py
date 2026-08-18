#!/usr/bin/env python3
"""Deduplicate manager deck: fewer slides, unique bullets only."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

DECK_PATH = Path(__file__).resolve().parent / "Document_Processing_POC_Manager_Deck_7_Slides.pptx"

FONT_BODY = Pt(20)
FONT_TITLE = Pt(28)


def _set_title(slide, title: str) -> None:
    sh = slide.shapes[0]
    tf = sh.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = FONT_TITLE


def _clear_extra_bullet_shapes(slide, keep_count: int) -> None:
    while len(slide.shapes) > 1 + keep_count:
        sp = slide.shapes[1]
        slide.shapes._spTree.remove(sp._element)


def _set_bullets(slide, bullets: list[str]) -> None:
    _clear_extra_bullet_shapes(slide, len(bullets))
    while len(slide.shapes) < 1 + len(bullets):
        from pptx.util import Emu

        idx = len(slide.shapes) - 1
        top = 1005840 + idx * 520000
        slide.shapes.add_textbox(731520, top, 10515600, 580000)
    for i, text in enumerate(bullets):
        sh = slide.shapes[1 + i]
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text if text.startswith("•") else f"• {text}"
        p.font.size = FONT_BODY


def _delete_slide(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    sld_id_lst.remove(ids[index])


def consolidate() -> None:
    prs = Presentation(str(DECK_PATH))

    # Slide 1 — Title
    _set_title(prs.slides[0], "Document Processing POC — Findings")
    _set_bullets(
        prs.slides[0],
        [
            "Confirm Excel, Word & PDF automation on an owned server + web app (not Matterway-only).",
            "Validate hosting path: API, logging, Windows/IIS deployment.",
        ],
    )

    # Slide 2 — Scope
    _set_title(prs.slides[1], "Objective & Scope")
    _set_bullets(
        prs.slides[1],
        [
            "Prove server-side Excel read/write and PDF/Word validation at scale.",
            "In scope: upload, server path, downloads, auto-download, in-place Excel, logs.",
            "Out of scope: full Matterway SCR scrub & CAL allocation products (port in next phase).",
        ],
    )

    # Slide 3 — Architecture (keep diagram shapes; only update first two bullet boxes)
    _set_title(prs.slides[2], "Solution Architecture")
    arch_bullets = [
        "React UI → FastAPI → Excel / PDF / Word services → results & downloads.",
        "REST API; config-driven PDF/Word rules; upload or authorized server path.",
    ]
    for i, text in enumerate(arch_bullets):
        sh = prs.slides[2].shapes[1 + i]
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"• {text}"
        p.font.size = FONT_BODY
    if len(prs.slides[2].shapes) > 3 and prs.slides[2].shapes[3].has_text_frame:
        prs.slides[2].shapes[3].text_frame.clear()

    # Slide 4 — Capabilities + outputs + findings (merge old 4,5,6)
    _set_title(prs.slides[3], "Capabilities & Deliverables")
    _set_bullets(
        prs.slides[3],
        [
            "Excel: POC Status, FALSE/NA highlights, Name/SSN tab; .xlsx/.xls; processed file kept for download.",
            "PDF/Word: PASS/FAIL and missing required phrases; metrics on result screen.",
            "Large files tested on server; upload copy removed after process; optional in-place workbook update.",
        ],
    )

    # Remove old slides: POC Findings (4), Business Outputs (5) — after merge, indices 4 & 5
    _delete_slide(prs, 5)
    _delete_slide(prs, 4)

    # Now index 4 = Matterway comparison
    _set_title(prs.slides[4], "POC vs Matterway (SCR & CAL)")
    _set_bullets(
        prs.slides[4],
        [
            "SCR: Helpwc scrubbing in-browser — POC proves server Excel; scrub rules not rebuilt here.",
            "CAL: multi-file allocation + YEDC/macro flows — POC proves PDF/Excel baseline only.",
            "POC adds: standalone app, API, server path, in-place update, enterprise hosting.",
        ],
    )

    # Merge Technology + Deployment (indices 5 & 6)
    _set_title(prs.slides[5], "Technology & Deployment")
    _set_bullets(
        prs.slides[5],
        [
            "Stack: React, FastAPI, openpyxl/xlrd, PyMuPDF, python-docx.",
            "Windows Server + IIS (HTTPS); FastAPI via NSSM on localhost:8000; reverse-proxy from IIS.",
            "Public port 443 only; runbook: docs/Windows-Server-Deployment-Guide.docx.",
        ],
    )
    _delete_slide(prs, 6)

    # Future (now index 6)
    _set_title(prs.slides[6], "Future Extensions")
    _set_bullets(
        prs.slides[6],
        [
            "Batch/scheduled jobs; notifications; auth (Azure AD / IIS).",
            "CAL-style multi-file wizards; PDF reports from Excel; phased SCR/CAL rule port + UAT.",
        ],
    )

    # Conclusion (index 7)
    _set_title(prs.slides[7], "Conclusion & Next Steps")
    _set_bullets(
        prs.slides[7],
        [
            "POC confirms document automation is feasible on our platform.",
            "Pilot one production workflow with auth, monitoring, and environment-based config.",
            "Migrate priority SCR/CAL logic onto this stack.",
        ],
    )

    prs.save(str(DECK_PATH))
    print(f"Consolidated {DECK_PATH} — {len(prs.slides)} slides.")


if __name__ == "__main__":
    consolidate()
