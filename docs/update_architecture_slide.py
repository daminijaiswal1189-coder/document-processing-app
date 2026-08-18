#!/usr/bin/env python3
"""Replace the simple architecture boxes with a clearer diagram image."""

from __future__ import annotations

import os
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
DECK_PATH = DOCS / "Document_Processing_POC_Manager_Deck_7_Slides.pptx"
DIAGRAM_PATH = DOCS / "architecture_diagram.png"
MPLCONFIGDIR = DOCS / ".matplotlib"

os.environ.setdefault("MPLCONFIGDIR", str(MPLCONFIGDIR))
MPLCONFIGDIR.mkdir(parents=True, exist_ok=True)


def _box(ax, x, y, w, h, text, fc, ec, fontsize=8.5, bold=False, text_color="#111827"):
    ax.add_patch(
        FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.015,rounding_size=0.07",
            linewidth=1.5,
            edgecolor=ec,
            facecolor=fc,
            zorder=2,
        )
    )
    ax.text(
        x + w / 2,
        y + h / 2,
        text,
        ha="center",
        va="center",
        fontsize=fontsize,
        fontweight="bold" if bold else "normal",
        color=text_color,
        linespacing=1.28,
        zorder=3,
    )


def _arrow(ax, x1, y1, x2, y2, color="#374151"):
    ax.add_patch(
        FancyArrowPatch(
            (x1, y1),
            (x2, y2),
            arrowstyle="-|>",
            mutation_scale=13,
            linewidth=1.35,
            color=color,
            zorder=1,
        )
    )


def _label(ax, x, y, text, color="#6B7280"):
    ax.text(x, y, text, fontsize=8.2, color=color, fontweight="bold", ha="left", va="center")


def render_diagram(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12.4, 6.0))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 6.0)
    ax.axis("off")

    # Soft layer backgrounds
    ax.add_patch(Rectangle((0.12, 4.55), 12.16, 1.25, facecolor="#F3F4F6", edgecolor="none", zorder=0))
    ax.add_patch(Rectangle((0.12, 2.55), 12.16, 1.85, facecolor="#EFF6FF", edgecolor="none", zorder=0))
    ax.add_patch(Rectangle((0.12, 0.12), 12.16, 2.28, facecolor="#ECFDF5", edgecolor="none", zorder=0))

    _label(ax, 0.25, 5.55, "1  Client & edge")
    _label(ax, 0.25, 4.15, "2  Application API", "#2563EB")
    _label(ax, 0.25, 2.15, "3  Document services & storage", "#059669")

    # --- Layer 1 ---
    _box(ax, 0.55, 4.75, 1.7, 0.7, "User\nBrowser", "#FFFFFF", "#4B5563", 9, True)
    _box(
        ax,
        3.0,
        4.75,
        3.2,
        0.7,
        "IIS (production)\nHTTPS + reverse proxy",
        "#FFEDD5",
        "#EA580C",
        8.8,
        True,
    )
    _box(
        ax,
        7.0,
        4.75,
        2.4,
        0.7,
        "React UI\nUpload · Path · Result",
        "#FEF3C7",
        "#D97706",
        8.8,
        True,
    )
    _box(
        ax,
        10.0,
        4.75,
        1.9,
        0.7,
        "Optional\nauto-download",
        "#FFFFFF",
        "#9CA3AF",
        8.2,
        True,
    )

    _arrow(ax, 2.25, 5.1, 3.0, 5.1)
    _arrow(ax, 6.2, 5.1, 7.0, 5.1)
    _arrow(ax, 9.4, 5.1, 10.0, 5.1)
    # IIS also fronts API
    ax.annotate(
        "",
        xy=(6.1, 3.95),
        xytext=(4.6, 4.75),
        arrowprops=dict(arrowstyle="-|>", color="#EA580C", lw=1.2),
    )
    ax.text(5.55, 4.45, "API routes", fontsize=7, color="#EA580C", ha="center")

    # --- Layer 2 ---
    _box(
        ax,
        2.4,
        3.05,
        7.6,
        0.95,
        "FastAPI (uvicorn)\n"
        "POST /upload · POST /upload/path · POST /process/excel|pdf|word\n"
        "GET /download/{file} · GET /health",
        "#DBEAFE",
        "#2563EB",
        8.6,
        True,
    )
    _arrow(ax, 8.2, 4.75, 6.8, 4.0)

    # --- Layer 3 processors ---
    _box(
        ax,
        0.45,
        1.35,
        2.55,
        0.85,
        "Excel processor\nopenpyxl · xlrd\nstatus · highlight · Name/SSN",
        "#DCFCE7",
        "#16A34A",
        7.8,
        True,
    )
    _box(
        ax,
        3.3,
        1.35,
        2.55,
        0.85,
        "PDF validator\nPyMuPDF\nPASS / FAIL + gaps",
        "#E0E7FF",
        "#4F46E5",
        7.8,
        True,
    )
    _box(
        ax,
        6.15,
        1.35,
        2.55,
        0.85,
        "Word validator\npython-docx\nPASS / FAIL + gaps",
        "#FCE7F3",
        "#DB2777",
        7.8,
        True,
    )
    _box(
        ax,
        9.1,
        1.35,
        2.7,
        0.85,
        "In-place Excel update\n(optional configured path)\nsame file on disk",
        "#F3F4F6",
        "#6B7280",
        7.6,
        True,
    )

    _arrow(ax, 4.0, 3.05, 1.7, 2.2)
    _arrow(ax, 5.6, 3.05, 4.55, 2.2)
    _arrow(ax, 7.0, 3.05, 7.4, 2.2)
    _arrow(ax, 8.5, 3.05, 10.3, 2.2)

    # Storage row — clear ownership
    _box(
        ax,
        0.55,
        0.25,
        2.9,
        0.75,
        "storage/uploads\nAPI stores file; cleaned after process",
        "#FFFFFF",
        "#059669",
        7.4,
    )
    _box(
        ax,
        3.8,
        0.25,
        3.0,
        0.75,
        "storage/processed\nExcel output kept for download",
        "#FFFFFF",
        "#059669",
        7.4,
    )
    _box(
        ax,
        7.15,
        0.25,
        2.5,
        0.75,
        "backend/logs\nTimestamped run logs",
        "#FFFFFF",
        "#059669",
        7.4,
    )
    _box(
        ax,
        10.0,
        0.25,
        1.85,
        0.75,
        "Config\nPDF/Word rules",
        "#FFFFFF",
        "#059669",
        7.4,
    )

    _arrow(ax, 3.5, 3.05, 2.0, 1.0)  # API → uploads
    _arrow(ax, 1.7, 1.35, 5.0, 1.0)  # Excel → processed
    _arrow(ax, 6.2, 3.05, 8.3, 1.0)  # API → logs
    _arrow(ax, 9.0, 3.05, 10.9, 1.0)  # API → config use

    ax.set_title(
        "Document Processing POC — end-to-end architecture",
        fontsize=12.5,
        fontweight="bold",
        color="#111827",
        pad=8,
    )
    fig.tight_layout(pad=0.25)
    fig.savefig(path, dpi=190, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def update_slide() -> None:
    render_diagram(DIAGRAM_PATH)
    prs = Presentation(str(DECK_PATH))
    slide = prs.slides[2]

    for shape in list(slide.shapes)[1:]:
        slide.shapes._spTree.remove(shape._element)

    caption = slide.shapes.add_textbox(Emu(550000), Emu(680000), Emu(11000000), Emu(380000))
    tf = caption.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Production: browser → IIS (HTTPS) → React UI + FastAPI. "
        "Each document type is handled by a dedicated service; Excel outputs stay available for download."
    )
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
    p.alignment = PP_ALIGN.LEFT

    slide.shapes.add_picture(
        str(DIAGRAM_PATH),
        Inches(0.28),
        Inches(1.25),
        width=Inches(12.15),
    )

    prs.save(str(DECK_PATH))
    print(f"Updated architecture slide in {DECK_PATH}")
    print(f"Diagram saved to {DIAGRAM_PATH}")


if __name__ == "__main__":
    update_slide()
